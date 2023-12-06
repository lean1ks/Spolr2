from pptx import Presentation

def replace_text_in_pptx(pptx_file_path, replacements):
    presentation = Presentation(pptx_file_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old_text, new_text in replacements.items():
                            run.text = run.text.replace(old_text, new_text)

    presentation.save('output.pptx')

# Пример использования
file_path = r'C:\Users\Администратор\Desktop\Шаблон сертификата.pptx'
text_replacements ={
    'FAM':'Фамилия',
    'Proffessia':'Заданая профессия'}

replace_text_in_pptx(file_path, text_replacements)
